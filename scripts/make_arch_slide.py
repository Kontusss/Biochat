"""Generate the one-page Biochat architecture & advantages slide (16:9 .pptx).

All shapes are native PowerPoint objects, so the deck stays fully editable.
Run: python scripts/make_arch_slide.py
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches, Pt

# ── palette ────────────────────────────────────────────────────────
INK        = RGBColor(0x0B, 0x2E, 0x3B)   # near-black navy, titles
INK_SOFT   = RGBColor(0x54, 0x6B, 0x76)   # secondary text
ACCENT     = RGBColor(0x0E, 0x9F, 0x8E)   # teal, the agent core
ACCENT_DK  = RGBColor(0x0A, 0x6B, 0x5F)
BLUE_DK    = RGBColor(0x1E, 0x4F, 0x6E)
VIOLET_DK  = RGBColor(0x4A, 0x43, 0x70)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
RULE       = RGBColor(0xD8, 0xE2, 0xE8)

LAYERS = {
    "ui":    (RGBColor(0xF1, 0xF6, 0xFA), RGBColor(0xCB, 0xDD, 0xE9), BLUE_DK),
    "svc":   (RGBColor(0xE8, 0xF1, 0xF7), RGBColor(0xBA, 0xD4, 0xE4), BLUE_DK),
    "agent": (RGBColor(0xE4, 0xF7, 0xF3), ACCENT,                      ACCENT_DK),
    "res":   (RGBColor(0xF3, 0xF2, 0xF8), RGBColor(0xCF, 0xCB, 0xE0), VIOLET_DK),
}

FONT = "PingFang SC"


# ── text helpers ───────────────────────────────────────────────────
def _style(run, size, bold=False, color=INK, font=FONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    rPr.get_or_add_latin().set("typeface", font)
    if rPr.find(qn("a:ea")) is None:
        rPr.find(qn("a:latin")).addnext(
            parse_xml(f'<a:ea {nsdecls("a")} typeface="{font}"/>')
        )


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color) or (text, size, bold, color, space_before)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[:4]
        space_before = spec[4] if len(spec) > 4 else 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(space_before)
        p.line_spacing = 1.18
        _style(p.add_run(), size, bold, color)
        p.runs[0].text = text
    return tb


def box(slide, x, y, w, h, fill, line, width=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        radius=0.055):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(width)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = radius
    s.text_frame.text = ""
    return s


_ARROW_SHAPE = {
    "r": MSO_SHAPE.RIGHT_ARROW, "l": MSO_SHAPE.LEFT_ARROW,
    "d": MSO_SHAPE.DOWN_ARROW,  "u": MSO_SHAPE.UP_ARROW,
}


def arrow(slide, x, y, w, h, direction, fill, label=None, label_y=None,
          label_color=INK_SOFT):
    s = slide.shapes.add_shape(_ARROW_SHAPE[direction],
                               Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    s.adjustments[0] = 0.42              # slim tail
    s.adjustments[1] = 0.30              # short head
    if label:
        textbox(slide, x - 0.15, label_y, w + 0.30, 0.22,
                [(label, 8.5, True, label_color)], align=PP_ALIGN.CENTER)
    return s


# ── deck ───────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# background
bg = box(slide, 0, 0, 13.333, 7.5, WHITE, WHITE, 0.25, MSO_SHAPE.RECTANGLE)
bg.shadow.inherit = False

# ── header ─────────────────────────────────────────────────────────
accent_bar = box(slide, 0.55, 0.42, 0.075, 0.62, ACCENT, ACCENT, 0.25,
                 MSO_SHAPE.RECTANGLE)
textbox(slide, 0.78, 0.38, 8.2, 0.72, [
    ("Biochat · 生物医学自主科研 Agent", 25, True, INK),
    ("规划 → 资源检索 → 代码执行 → 观察反思 的闭环架构", 11.5, False, INK_SOFT, 4),
])
textbox(slide, 9.4, 0.46, 3.4, 0.6, [
    ("226 工具 / 23 领域 · 76 数据集 · 113 软件库", 10, True, ACCENT_DK),
    ("8 家 LLM 供应商可插拔 · MCP 协议接入", 10, False, INK_SOFT, 3),
], align=PP_ALIGN.RIGHT)

rule = box(slide, 0.55, 1.22, 12.23, 0.012, RULE, RULE, 0.25, MSO_SHAPE.RECTANGLE)

# ═══ LEFT: architecture diagram ════════════════════════════════════
DX, DW = 0.55, 7.15          # diagram x / width
PAD = 0.22                   # inner text padding

# — layer 1: 交互层 —
y = 1.46
box(slide, DX, y, DW, 0.66, *LAYERS["ui"][:2])
textbox(slide, DX + PAD, y + 0.13, 1.5, 0.4, [("交互层", 12, True, LAYERS["ui"][2])])
textbox(slide, DX + 1.55, y + 0.15, DW - 1.75, 0.4,
        [("Streamlit UI（流式对话 / 多会话 / PDF 导出） · Gradio · MCP Server", 10.5, False, INK)])

# — layer 2: 服务层 —
y = 2.36
box(slide, DX, y, DW, 0.66, *LAYERS["svc"][:2])
textbox(slide, DX + PAD, y + 0.13, 1.5, 0.4, [("服务层", 12, True, LAYERS["svc"][2])])
textbox(slide, DX + 1.55, y + 0.15, DW - 1.75, 0.4,
        [("BioAgentService — 生命周期 · 增量流式事件 · 出口净化（内部推理不外泄）", 10.5, False, INK)])

# — layer 3: Agent 引擎（核心，加重描边）—
AY, AH = 3.26, 2.52
box(slide, DX, AY, DW, AH, LAYERS["agent"][0], LAYERS["agent"][1], 1.75)
textbox(slide, DX + PAD, AY + 0.14, 5.6, 0.4,
        [("Agent 引擎 · A1（LangGraph 状态机）", 12, True, ACCENT_DK)])
textbox(slide, DX + DW - 2.5, AY + 0.16, 2.3, 0.4,
        [("单步超时 / 解析重试保护", 8.5, False, INK_SOFT)], align=PP_ALIGN.RIGHT)

# generate ↔ execute loop
NY, NH, NW = 3.82, 1.00, 2.25
gen = box(slide, DX + 0.42, NY, NW, NH, WHITE, ACCENT, 1.4)
textbox(slide, DX + 0.42, NY + 0.20, NW, 0.62, [
    ("generate", 13, True, ACCENT_DK),
    ("LLM 规划 · 生成代码", 9, False, INK_SOFT, 3),
], align=PP_ALIGN.CENTER)

exe = box(slide, DX + DW - 0.42 - NW, NY, NW, NH, WHITE, ACCENT, 1.4)
textbox(slide, DX + DW - 0.42 - NW, NY + 0.20, NW, 0.62, [
    ("execute", 13, True, ACCENT_DK),
    ("Python / R / Bash 沙箱", 9, False, INK_SOFT, 3),
], align=PP_ALIGN.CENTER)

gap_x = DX + 0.42 + NW
gap_w = (DX + DW - 0.42 - NW) - gap_x
arrow(slide, gap_x + 0.14, NY + 0.24, gap_w - 0.28, 0.20, "r", ACCENT,
      "<execute>", NY + 0.02)
arrow(slide, gap_x + 0.14, NY + 0.66, gap_w - 0.28, 0.20, "l",
      RGBColor(0x7F, 0xB8, 0xB0), "<observation>", NY + 0.86)

# solution exit
arrow(slide, DX + 0.42 + NW / 2 - 0.11, NY + NH + 0.03, 0.22, 0.30, "d", ACCENT)
sol = box(slide, DX + 0.42, NY + NH + 0.36, NW, 0.40, ACCENT, ACCENT, 0.75)
textbox(slide, DX + 0.42, NY + NH + 0.46, NW, 0.26,
        [("<solution> → 结构化交付", 10.5, True, WHITE)], align=PP_ALIGN.CENTER)

textbox(slide, gap_x + 0.30, NY + NH + 0.40, DW - NW - 0.9, 0.42, [
    ("六段式输出：结论 / 依据与原理 / 方法摘要 /", 9.5, False, INK_SOFT),
    ("建议验证实验 / 不确定性与局限 / 安全声明", 9.5, False, INK_SOFT, 2),
])

# — layer 4: 资源层 —
y = 6.10
box(slide, DX, y, DW, 0.92, *LAYERS["res"][:2])
textbox(slide, DX + PAD, y + 0.11, 4.0, 0.3,
        [("资源层 · ResourceSelector 按需检索", 12, True, VIOLET_DK)])
textbox(slide, DX + PAD, y + 0.48, DW - 2 * PAD, 0.3,
        [("226 个工具函数（23 领域）   ·   76 个生物医学数据集   ·   113 个软件库   ·   Know-how 知识库", 10, False, INK)])

# feed arrow: 资源层 → Agent 引擎（在两层之间的 0.32" 间隙内）
arrow(slide, DX + DW - 1.60, AY + AH + 0.04, 0.20, 0.24, "u",
      RGBColor(0xA9, 0xA2, 0xC4))
textbox(slide, DX + DW - 3.85, AY + AH + 0.08, 2.10, 0.22,
        [("仅注入相关子集", 8.5, True, VIOLET_DK)], align=PP_ALIGN.RIGHT)

# ═══ RIGHT: differentiated advantages ══════════════════════════════
RX, RW = 8.12, 4.66
textbox(slide, RX, 1.46, RW, 0.3, [("核心优势", 13, True, INK)])
box(slide, RX, 1.80, 0.42, 0.028, ACCENT, ACCENT, 0.25, MSO_SHAPE.RECTANGLE)

CARDS = [
    ("代码级执行，而非「生成答案」", [
        "每个结论都由真实运行的 Python/R 产出，observation 回灌",
        "下一轮推理 — 过程可复现、可审计，不靠模型记忆作答。"]),
    ("按需检索，不把工具全塞进上下文", [
        "ResourceSelector 一次 LLM 调用即从 226 工具 / 76 数据集中",
        "筛出相关子集 — 省 token、降幻觉、工具库可横向扩容。"]),
    ("垂直管线落地，一句话跑完全流程", [
        "抗体设计：CDRH3 扩散生成 → 结构预测 → HDock 对接 →",
        "可开发性评估 → 多维排序，非通用 LLM 可替代的深度。"]),
    ("工程可控，可直接进生产", [
        "出口净化器杜绝内部推理泄漏 · 禁止编造 Kd/ΔG 等指标 ·",
        "8 家 LLM 供应商热插拔 · 商业/非商业数据自动隔离。"]),
]

cy = 2.02
for i, (title, body_lines) in enumerate(CARDS, 1):
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(RX), Inches(cy + 0.02),
                                   Inches(0.30), Inches(0.30))
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    badge.shadow.inherit = False
    btf = badge.text_frame
    btf.margin_left = btf.margin_right = btf.margin_top = btf.margin_bottom = 0
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _style(btf.paragraphs[0].add_run(), 10.5, True, WHITE)
    btf.paragraphs[0].runs[0].text = str(i)

    textbox(slide, RX + 0.44, cy, RW - 0.44, 0.9,
            [(title, 12.5, True, INK)]
            + [(ln, 9.5, False, INK_SOFT, 5 if j == 0 else 0)
               for j, ln in enumerate(body_lines)])
    if i < len(CARDS):
        box(slide, RX + 0.44, cy + 1.13, RW - 0.44, 0.012, RULE, RULE, 0.25,
            MSO_SHAPE.RECTANGLE)
    cy += 1.33

# ── footer ─────────────────────────────────────────────────────────
box(slide, 0.55, 7.16, 12.23, 0.012, RULE, RULE, 0.25, MSO_SHAPE.RECTANGLE)
textbox(slide, 0.55, 7.24, 8.5, 0.24,
        [("技术栈：LangGraph · LangChain · Pydantic · Streamlit　|　工程化：ruff + pre-commit + pytest", 8.5, False, INK_SOFT)])
textbox(slide, 9.0, 7.24, 3.78, 0.24,
        [("基于 Biomni (Apache 2.0) 深度重构", 8.5, False, INK_SOFT)], align=PP_ALIGN.RIGHT)

OUT = "Biochat_架构与优势.pptx"
prs.save(OUT)
print(f"saved → {OUT}")
