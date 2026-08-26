"""Shared python-pptx helpers for the slide generators in this directory.

All shapes produced here are native PowerPoint objects, so the decks stay
fully editable. Both the latin and east-asian typefaces are set on every run,
which is what keeps CJK text from falling back to a substitute font.
"""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Inches, Pt

# ── palette ────────────────────────────────────────────────────────
INK = RGBColor(0x0B, 0x2E, 0x3B)
INK_SOFT = RGBColor(0x54, 0x6B, 0x76)
ACCENT = RGBColor(0x0E, 0x9F, 0x8E)
ACCENT_DK = RGBColor(0x0A, 0x6B, 0x5F)
BLUE_DK = RGBColor(0x1E, 0x4F, 0x6E)
VIOLET_DK = RGBColor(0x4A, 0x43, 0x70)
AMBER_DK = RGBColor(0x9A, 0x5B, 0x0E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xD8, 0xE2, 0xE8)

FONT = "PingFang SC"

_ARROW_SHAPE = {
    "r": MSO_SHAPE.RIGHT_ARROW,
    "l": MSO_SHAPE.LEFT_ARROW,
    "d": MSO_SHAPE.DOWN_ARROW,
    "u": MSO_SHAPE.UP_ARROW,
}


def style_run(run, size, bold=False, color=INK, font=FONT):
    """Apply size/weight/colour and pin both latin and east-asian typefaces."""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    rPr.get_or_add_latin().set("typeface", font)
    if rPr.find(qn("a:ea")) is None:
        rPr.find(qn("a:latin")).addnext(parse_xml(f'<a:ea {nsdecls("a")} typeface="{font}"/>'))


def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.18):
    """lines: sequence of ``(text, size, bold, color[, space_before_pt])``."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, spec in enumerate(lines):
        text, size, bold, color = spec[:4]
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(spec[4] if len(spec) > 4 else 0)
        p.line_spacing = spacing
        style_run(p.add_run(), size, bold, color)
        p.runs[0].text = text
    return tb


def box(slide, x, y, w, h, fill, line, width=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.055):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(width)
    s.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        s.adjustments[0] = radius
    s.text_frame.text = ""
    return s


def arrow(slide, x, y, w, h, direction, fill, label=None, label_y=None, label_color=INK_SOFT):
    s = slide.shapes.add_shape(_ARROW_SHAPE[direction], Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    s.adjustments[0] = 0.42
    s.adjustments[1] = 0.30
    if label:
        textbox(slide, x - 0.15, label_y, w + 0.30, 0.22, [(label, 8.5, True, label_color)], align=PP_ALIGN.CENTER)
    return s


def badge(slide, x, y, diameter, text, fill=ACCENT, color=WHITE, size=10.5):
    """Small filled circle with centred text — used for step numbers."""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    tf = s.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    style_run(tf.paragraphs[0].add_run(), size, True, color)
    tf.paragraphs[0].runs[0].text = text
    return s


def rule(slide, x, y, w, color=RULE):
    return box(slide, x, y, w, 0.012, color, color, 0.25, MSO_SHAPE.RECTANGLE)


def new_deck():
    """A blank 16:9 presentation with one empty slide and a white background."""
    from pptx import Presentation

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box(slide, 0, 0, 13.333, 7.5, WHITE, WHITE, 0.25, MSO_SHAPE.RECTANGLE)
    return prs, slide
